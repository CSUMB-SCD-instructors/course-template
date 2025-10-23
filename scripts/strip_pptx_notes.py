#!/usr/bin/env python3
"""
Strip speaker notes from PowerPoint presentations (.pptx files).

This script removes all speaker notes from PowerPoint files, allowing you to
share slides with students while keeping instructor notes private.

Usage:
    python strip_ppt_notes.py input.pptx                    # Creates input_no_notes.pptx
    python strip_ppt_notes.py input.pptx -o output.pptx    # Specify output file
    python strip_ppt_notes.py slides/*.pptx                # Process multiple files
    python strip_ppt_notes.py input.pptx --in-place        # Overwrite original file

Requirements:
    pip install python-pptx
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx library not found.", file=sys.stderr)
    print("Install it with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def strip_notes(input_path: Path, output_path: Path) -> bool:
    """
    Strip all speaker notes from a PowerPoint file.

    Args:
        input_path: Path to input PPTX file
        output_path: Path to output PPTX file

    Returns:
        True if notes were found and removed, False if no notes existed
    """
    try:
        prs = Presentation(input_path)
        notes_found = False

        for slide in prs.slides:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame

                # Check if there's actual content in the notes
                if text_frame.text.strip():
                    notes_found = True

                # Clear all paragraphs in the notes
                text_frame.clear()

        prs.save(output_path)
        return notes_found

    except Exception as e:
        print(f"Error processing {input_path}: {e}", file=sys.stderr)
        return False


def main():
    parser = argparse.ArgumentParser(
        description="Strip speaker notes from PowerPoint files",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  Strip notes and create new file:
    %(prog)s lecture01.pptx

  Specify output filename:
    %(prog)s lecture01.pptx -o lecture01_student.pptx

  Process multiple files:
    %(prog)s slides/*.pptx

  Overwrite original file (use with caution):
    %(prog)s lecture01.pptx --in-place
        """
    )

    parser.add_argument(
        'input_files',
        nargs='+',
        type=Path,
        help='PowerPoint file(s) to process (.pptx)'
    )

    parser.add_argument(
        '-o', '--output',
        type=Path,
        help='Output file path (only valid with single input file)'
    )

    parser.add_argument(
        '--in-place',
        action='store_true',
        help='Overwrite original file(s) instead of creating new files'
    )

    parser.add_argument(
        '-q', '--quiet',
        action='store_true',
        help='Suppress output messages'
    )

    args = parser.parse_args()

    # Validate arguments
    if args.output and len(args.input_files) > 1:
        parser.error("Cannot specify --output with multiple input files")

    if args.output and args.in_place:
        parser.error("Cannot use both --output and --in-place")

    # Process each file
    success_count = 0
    error_count = 0

    for input_path in args.input_files:
        if not input_path.exists():
            print(f"Error: File not found: {input_path}", file=sys.stderr)
            error_count += 1
            continue

        if not input_path.suffix.lower() in ['.pptx', '.ppt']:
            print(f"Warning: Skipping {input_path} (not a .pptx file)", file=sys.stderr)
            continue

        if input_path.suffix.lower() == '.ppt':
            print(f"Error: Old PowerPoint format (.ppt) not supported. Please convert {input_path} to .pptx first.", file=sys.stderr)
            error_count += 1
            continue

        # Determine output path
        if args.in_place:
            output_path = input_path
        elif args.output:
            output_path = args.output
        else:
            # Default: add _no_notes before extension
            stem = input_path.stem
            output_path = input_path.parent / f"{stem}_no_notes{input_path.suffix}"

        # Process the file
        notes_found = strip_notes(input_path, output_path)

        if notes_found is not False:  # Success (either notes found or not)
            success_count += 1
            if not args.quiet:
                if notes_found:
                    print(f"✓ Stripped notes from {input_path} → {output_path}")
                else:
                    print(f"✓ No notes found in {input_path} → {output_path}")
        else:  # Error
            error_count += 1

    # Summary
    if not args.quiet and len(args.input_files) > 1:
        print(f"\nProcessed {success_count} file(s) successfully, {error_count} error(s)")

    sys.exit(0 if error_count == 0 else 1)


if __name__ == "__main__":
    main()